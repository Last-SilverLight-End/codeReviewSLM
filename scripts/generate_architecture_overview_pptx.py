from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/architecture_overview.pptx")
FONT = "Malgun Gothic"

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(246, 249, 253)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(12, 35, 70)
TEXT = RGBColor(27, 38, 56)
MUTED = RGBColor(93, 108, 132)
LINE = RGBColor(185, 199, 219)
SOFT = RGBColor(235, 241, 249)

BLUE = RGBColor(42, 104, 190)
BLUE_SOFT = RGBColor(226, 238, 253)
GREEN = RGBColor(38, 150, 112)
GREEN_SOFT = RGBColor(226, 246, 239)
PURPLE = RGBColor(120, 91, 190)
PURPLE_SOFT = RGBColor(238, 233, 250)
YELLOW = RGBColor(214, 160, 43)
YELLOW_SOFT = RGBColor(255, 248, 225)
TEAL = RGBColor(37, 143, 168)
TEAL_SOFT = RGBColor(225, 247, 250)
ORANGE = RGBColor(210, 122, 45)
ORANGE_SOFT = RGBColor(252, 238, 223)
RED = RGBColor(202, 77, 77)


def apply_font(run, size=12, color=TEXT, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_text(slide, x, y, w, h, value, size=12, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    apply_font(p.runs[0], size=size, color=color, bold=bold)
    return shape


def add_box(
    slide,
    x,
    y,
    w,
    h,
    value,
    fill=WHITE,
    border=LINE,
    size=12,
    color=TEXT,
    bold=True,
    align=PP_ALIGN.CENTER,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.15)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    apply_font(p.runs[0], size=size, color=color, bold=bold)
    return shape


def add_panel(slide, x, y, w, h, title=None, fill=WHITE, border=LINE, title_color=NAVY):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = border
    panel.line.width = Pt(1.1)
    if title:
        add_text(slide, x + 0.16, y + 0.08, w - 0.32, 0.26, title, 11, title_color, True, PP_ALIGN.LEFT)
    return panel


def add_bullets(slide, x, y, w, h, items, size=9.8, color=TEXT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.space_after = Pt(1)
        apply_font(p.runs[0], size=size, color=color, bold=False)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.5):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    connector.line.end_arrowhead = True
    return connector


def add_title(slide, title, subtitle):
    add_text(slide, 0.55, 0.16, 12.25, 0.42, title, 28, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, 0.8, 0.62, 11.75, 0.28, subtitle, 12.5, MUTED, False, PP_ALIGN.CENTER)


def add_footer(slide, number):
    add_text(slide, 0.6, 7.15, 12.1, 0.18, f"Local AI Code Review Platform Architecture | {number}", 8, MUTED, False, PP_ALIGN.RIGHT)


def overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Local AI Code Review Platform", "로컬 LLM 기반 코드 리뷰 · Hybrid RAG · 운영 관측성")

    # Outer main area
    add_panel(slide, 0.35, 1.02, 12.65, 5.95, fill=WHITE, border=LINE)

    # Top flow
    add_text(slide, 0.65, 1.18, 1.1, 0.24, "사용자 흐름", 10.5, NAVY, True)
    flow_y = 1.48
    steps = [
        ("사용자", BLUE_SOFT, BLUE),
        ("Next.js UI\n업로드/채팅", BLUE_SOFT, BLUE),
        ("FastAPI\nAPI Gateway", GREEN_SOFT, GREEN),
        ("Hybrid\nRetriever", PURPLE_SOFT, PURPLE),
        ("Ollama\nLLM 답변", ORANGE_SOFT, ORANGE),
        ("Kibana\n관측", YELLOW_SOFT, YELLOW),
    ]
    step_x = [0.85, 2.75, 4.95, 7.15, 9.25, 11.05]
    step_w = [1.25, 1.5, 1.55, 1.45, 1.45, 1.25]
    for i, (label, fill, border) in enumerate(steps):
        add_box(slide, step_x[i], flow_y, step_w[i], 0.65, label, fill, border, 10.5, NAVY)
        if i < len(steps) - 1:
            add_arrow(slide, step_x[i] + step_w[i] + 0.08, flow_y + 0.33, step_x[i + 1] - 0.08, flow_y + 0.33, NAVY, 1.25)

    # Left operation lane
    add_panel(slide, 0.65, 2.55, 2.15, 3.0, "작업 실행", fill=SOFT, border=LINE)
    add_box(slide, 0.9, 3.02, 1.65, 0.45, "인증 / 세션", WHITE, BLUE, 9.8, NAVY)
    add_box(slide, 0.9, 3.62, 1.65, 0.45, "파일 / ZIP 업로드", WHITE, BLUE, 9.8, NAVY)
    add_box(slide, 0.9, 4.22, 1.65, 0.45, "비동기 리뷰", WHITE, BLUE, 9.8, NAVY)
    add_box(slide, 0.9, 4.82, 1.65, 0.45, "SSE 스트리밍", WHITE, BLUE, 9.8, NAVY)

    # Center analysis lane
    add_panel(slide, 3.05, 2.55, 5.95, 3.0, "AI 코드 분석 핵심", fill=WHITE, border=LINE)
    add_box(slide, 3.35, 3.02, 1.55, 0.52, "Tree-sitter\n청킹", TEAL_SOFT, TEAL, 9.5, NAVY)
    add_box(slide, 5.15, 3.02, 1.55, 0.52, "Ollama\n임베딩", ORANGE_SOFT, ORANGE, 9.5, NAVY)
    add_box(slide, 6.95, 3.02, 1.55, 0.52, "Context\nBuilder", GREEN_SOFT, GREEN, 9.5, NAVY)
    add_arrow(slide, 4.9, 3.28, 5.15, 3.28, NAVY, 1.15)
    add_arrow(slide, 6.7, 3.28, 6.95, 3.28, NAVY, 1.15)

    add_box(slide, 3.55, 4.12, 1.75, 0.5, "pgvector\n의미 검색", PURPLE_SOFT, PURPLE, 9.5, NAVY)
    add_box(slide, 5.62, 4.12, 1.75, 0.5, "Elasticsearch\nBM25 검색", YELLOW_SOFT, YELLOW, 9.5, NAVY)
    add_box(slide, 7.72, 4.12, 0.95, 0.5, "RRF\n병합", GREEN_SOFT, GREEN, 9.5, NAVY)
    add_arrow(slide, 5.3, 4.37, 5.62, 4.37, NAVY, 1.15)
    add_arrow(slide, 7.37, 4.37, 7.72, 4.37, NAVY, 1.15)
    add_text(slide, 3.55, 4.86, 5.0, 0.2, "검색 가중치: Elasticsearch 1.25 / pgvector 1.0", 8.8, MUTED, False, PP_ALIGN.CENTER)

    # Right observability lane
    add_panel(slide, 9.28, 2.55, 3.25, 3.0, "저장소 / 관측", fill=SOFT, border=LINE)
    add_box(slide, 9.58, 3.02, 1.18, 0.5, "PostgreSQL\n원본/권한", PURPLE_SOFT, PURPLE, 8.8, NAVY)
    add_box(slide, 10.98, 3.02, 1.18, 0.5, "Redis\n세션/브로커", WHITE, LINE, 8.8, NAVY)
    add_box(slide, 9.58, 3.82, 1.18, 0.5, "Elasticsearch\n검색/로그", YELLOW_SOFT, YELLOW, 8.8, NAVY)
    add_box(slide, 10.98, 3.82, 1.18, 0.5, "Kibana\n대시보드", YELLOW_SOFT, YELLOW, 8.8, NAVY)
    add_box(slide, 9.58, 4.62, 2.58, 0.5, "codereview-* / code-chunks", WHITE, LINE, 9.2, NAVY)

    # Cross arrows
    add_arrow(slide, 2.8, 4.05, 3.05, 4.05, BLUE, 1.15)
    add_arrow(slide, 9.0, 4.05, 9.28, 4.05, YELLOW, 1.15)
    add_arrow(slide, 8.05, 2.13, 8.05, 2.55, GREEN, 1.0)

    # Bottom principles
    add_panel(slide, 0.65, 5.85, 11.88, 0.72, "설계 원칙", fill=WHITE, border=LINE)
    principles = [
        ("로컬 실행", "Ollama Native"),
        ("Hybrid RAG", "의미 + 키워드"),
        ("비동기 처리", "Celery Worker"),
        ("관측 가능", "Elasticsearch/Kibana"),
        ("권한 기준", "PostgreSQL"),
        ("실시간 응답", "SSE"),
    ]
    x = 1.0
    for title, sub in principles:
        add_box(slide, x, 6.2, 1.6, 0.26, title, WHITE, LINE, 8.2, NAVY)
        add_text(slide, x, 6.43, 1.6, 0.13, sub, 6.8, MUTED, False, PP_ALIGN.CENTER)
        x += 1.9

    add_footer(slide, 1)
    return slide


def detail_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "아키텍처 해석", "각 계층이 담당하는 책임과 데이터 흐름")

    add_panel(slide, 0.65, 1.18, 12.05, 5.72, fill=WHITE, border=LINE)
    rows = [
        ("Frontend", "Next.js UI", "사용자 입력, 파일/이미지 업로드, SSE 응답 표시, RAG 참조 표시"),
        ("Backend", "FastAPI", "인증, 코드 파싱, 대화/리뷰 API, Hybrid Retriever 실행"),
        ("Vector", "PostgreSQL + pgvector", "원본 데이터와 임베딩 저장, 의미 기반 코드 검색"),
        ("Keyword", "Elasticsearch", "코드 청크 BM25 검색, API/RAG/LLM 로그 저장"),
        ("Model", "Ollama", "qwen3 답변 생성, nomic embedding, llava 이미지 분석"),
        ("Async", "Celery + Redis", "긴 코드 리뷰 작업 분리, refresh token 및 broker 처리"),
        ("Observe", "Kibana", "인덱스 조회, 실패율, 지연 시간, RAG 결과 수 관측"),
    ]
    y = 1.65
    for key, name, desc in rows:
        add_box(slide, 0.95, y, 1.35, 0.42, key, BLUE_SOFT, BLUE, 9, NAVY)
        add_text(slide, 2.55, y, 2.25, 0.42, name, 10, NAVY, True)
        add_text(slide, 4.85, y, 7.35, 0.42, desc, 9.2, TEXT)
        y += 0.67

    add_footer(slide, 2)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    overview_slide(prs)
    detail_slide(prs)
    return prs


if __name__ == "__main__":
    presentation = build()
    presentation.save(OUT)
    print(OUT)

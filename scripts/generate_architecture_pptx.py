from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/architecture.pptx")
FONT = "Malgun Gothic"

BG = RGBColor(12, 16, 23)
PANEL = RGBColor(29, 36, 49)
PANEL_DARK = RGBColor(21, 26, 36)
LINE = RGBColor(78, 92, 120)
TEXT = RGBColor(242, 245, 250)
MUTED = RGBColor(170, 180, 198)
BLUE = RGBColor(67, 135, 245)
GREEN = RGBColor(67, 170, 112)
YELLOW = RGBColor(218, 174, 73)
PURPLE = RGBColor(148, 112, 219)
CYAN = RGBColor(70, 178, 190)
ORANGE = RGBColor(220, 130, 68)


def set_font(run, size=None, color=None, bold=False):
    run.font.name = FONT
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    run.font.bold = bold


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def text_box(slide, x, y, w, h, text, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_font(p.runs[0], size=size, color=color, bold=bold)
    return shape


def title(slide, main, sub=None):
    text_box(slide, 0.55, 0.28, 12.3, 0.5, main, 26, TEXT, True)
    if sub:
        text_box(slide, 0.57, 0.82, 12.1, 0.34, sub, 11, MUTED)


def footer(slide, number):
    text_box(
        slide,
        0.65,
        7.06,
        12.0,
        0.22,
        f"Local AI Code Review Platform Architecture  |  {number}",
        8,
        RGBColor(104, 116, 138),
        False,
        PP_ALIGN.RIGHT,
    )


def box(slide, x, y, w, h, text, fill=PANEL, border=LINE, size=12, bold=True, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_font(p.runs[0], size=size, color=TEXT, bold=bold)
    return shape


def bullet_panel(slide, x, y, w, h, head, items, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.3)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = head
    set_font(p.runs[0], size=14, color=TEXT, bold=True)
    for item in items:
        bp = tf.add_paragraph()
        bp.text = f"• {item}"
        set_font(bp.runs[0], size=10.5, color=MUTED)
    return shape


def arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.6):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    connector.line.end_arrowhead = True
    return connector


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text_box(slide, 0.8, 1.1, 11.8, 0.35, "로컬 AI 코드 리뷰 플랫폼", 18, BLUE, True, PP_ALIGN.CENTER)
    text_box(slide, 0.8, 1.65, 11.8, 0.8, "전체 아키텍처", 42, TEXT, True, PP_ALIGN.CENTER)
    text_box(
        slide,
        1.6,
        2.65,
        10.1,
        0.45,
        "Next.js · FastAPI · Ollama · pgvector · Elasticsearch · Kibana",
        17,
        MUTED,
        False,
        PP_ALIGN.CENTER,
    )
    box(slide, 2.15, 4.05, 1.75, 0.68, "UI", BLUE, BLUE, 14)
    box(slide, 4.25, 4.05, 1.75, 0.68, "API", GREEN, GREEN, 14)
    box(slide, 6.35, 4.05, 1.75, 0.68, "RAG", PURPLE, PURPLE, 14)
    box(slide, 8.45, 4.05, 1.75, 0.68, "Logs", YELLOW, YELLOW, 14)
    arrow(slide, 3.9, 4.39, 4.25, 4.39)
    arrow(slide, 6.0, 4.39, 6.35, 4.39)
    arrow(slide, 8.1, 4.39, 8.45, 4.39)
    footer(slide, 1)

    # 2. 전체 아키텍처
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "전체 아키텍처 한눈에 보기", "사용자 요청이 UI, API, 검색, 모델, 로그 계층을 거쳐 답변과 관측 데이터로 돌아오는 구조")
    box(slide, 0.45, 1.45, 1.55, 0.65, "사용자", PANEL_DARK, LINE, 12)
    box(slide, 2.35, 1.25, 2.05, 1.05, "Next.js UI\n:3000\n채팅 · 업로드 · 참조", BLUE, BLUE, 11)
    arrow(slide, 2.0, 1.78, 2.35, 1.78, BLUE)

    backend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.95), Inches(0.95), Inches(3.25), Inches(5.75))
    backend.fill.solid()
    backend.fill.fore_color.rgb = RGBColor(23, 31, 42)
    backend.line.color.rgb = GREEN
    backend.line.width = Pt(1.5)
    text_box(slide, 5.1, 1.05, 2.95, 0.35, "FastAPI Backend :8000", 14, TEXT, True, PP_ALIGN.CENTER)
    box(slide, 5.25, 1.6, 1.15, 0.52, "Auth", GREEN, GREEN, 10)
    box(slide, 6.65, 1.6, 1.15, 0.52, "Code API", GREEN, GREEN, 10)
    box(slide, 5.25, 2.35, 1.15, 0.52, "Chat/RAG", GREEN, GREEN, 10)
    box(slide, 6.65, 2.35, 1.15, 0.52, "Review", GREEN, GREEN, 10)
    box(slide, 5.25, 3.1, 2.55, 0.55, "Hybrid Retriever", PURPLE, PURPLE, 11)
    box(slide, 5.25, 3.9, 1.15, 0.52, "SSE", CYAN, CYAN, 10)
    box(slide, 6.65, 3.9, 1.15, 0.52, "Admin Log", YELLOW, YELLOW, 10)
    box(slide, 5.25, 4.65, 2.55, 0.55, "Celery Worker", ORANGE, ORANGE, 11)
    arrow(slide, 4.4, 1.78, 4.95, 1.78, BLUE)

    box(slide, 8.9, 1.05, 2.0, 0.65, "Ollama :11434\nqwen3 · llava · embedding", PURPLE, PURPLE, 10)
    box(slide, 8.9, 2.05, 2.0, 0.65, "PostgreSQL + pgvector\n:5433", PURPLE, PURPLE, 10)
    box(slide, 8.9, 3.05, 2.0, 0.65, "Elasticsearch\n:9200", YELLOW, YELLOW, 10)
    box(slide, 8.9, 4.05, 2.0, 0.65, "Redis\n:6379", PANEL_DARK, LINE, 10)
    box(slide, 8.9, 5.05, 2.0, 0.65, "Kibana\n:5601", YELLOW, YELLOW, 10)
    bullet_panel(slide, 11.2, 1.05, 1.75, 1.65, "저장", ["원본 데이터", "임베딩", "대화/리뷰"], PURPLE)
    bullet_panel(slide, 11.2, 3.0, 1.75, 1.65, "검색/로그", ["BM25", "API 로그", "LLM 로그"], YELLOW)
    bullet_panel(slide, 11.2, 5.0, 1.75, 1.25, "답변", ["코드 기반", "RAG 참조"], BLUE)
    arrow(slide, 8.2, 1.45, 8.9, 1.38, PURPLE)
    arrow(slide, 8.2, 3.35, 8.9, 2.38, PURPLE)
    arrow(slide, 8.2, 3.35, 8.9, 3.38, YELLOW)
    arrow(slide, 8.2, 4.92, 8.9, 4.38, ORANGE)
    arrow(slide, 10.9, 5.38, 11.2, 5.6, YELLOW)

    box(slide, 0.65, 6.05, 2.15, 0.45, "파일/ZIP 업로드", PANEL_DARK, LINE, 9)
    box(slide, 3.0, 6.05, 2.15, 0.45, "청킹 + 임베딩", PANEL_DARK, LINE, 9)
    box(slide, 5.35, 6.05, 2.15, 0.45, "Hybrid 검색", PANEL_DARK, LINE, 9)
    box(slide, 7.7, 6.05, 2.15, 0.45, "LLM 답변", PANEL_DARK, LINE, 9)
    box(slide, 10.05, 6.05, 2.15, 0.45, "Kibana 관측", PANEL_DARK, LINE, 9)
    for x in [2.8, 5.15, 7.5, 9.85]:
        arrow(slide, x, 6.27, x + 0.2, 6.27, MUTED, 1.2)
    footer(slide, 2)

    # 3. 실행 경계
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "실행 경계", "Docker Compose와 Windows Native 프로세스를 분리")
    bullet_panel(slide, 1.0, 1.35, 5.2, 4.7, "Docker Compose", ["PostgreSQL + pgvector :5433", "Redis :6379", "Elasticsearch :9200", "Kibana :5601"], GREEN)
    bullet_panel(slide, 7.1, 1.35, 5.2, 4.7, "Windows Native", ["Ollama :11434", "FastAPI :8000", "Celery Worker", "Next.js :3000"], BLUE)
    text_box(slide, 1.1, 6.35, 11.1, 0.35, "Ollama는 Docker 컨테이너가 아니며, 컨테이너 내부 접근 시 host.docker.internal:11434 사용", 13, TEXT, True, PP_ALIGN.CENTER)
    footer(slide, 3)

    # 4. 데이터 저장 역할
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "데이터 저장 역할 분리")
    bullet_panel(slide, 0.8, 1.35, 5.5, 4.9, "PostgreSQL + pgvector", ["원본 데이터 기준 저장소", "users / projects / files / chunks", "reviews / conversations / messages", "사용자/프로젝트 권한 필터", "임베딩 기반 의미 검색"], PURPLE)
    bullet_panel(slide, 7.0, 1.35, 5.5, 4.9, "Elasticsearch + Kibana", ["codereview-code-chunks BM25 검색", "함수명 / 파일명 / 에러 문자열 검색", "API / RAG / LLM / WebSearch 로그", "Kibana data view", "운영 관측성"], YELLOW)
    footer(slide, 4)

    # 5. Hybrid Retrieval
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Hybrid Retrieval", "pgvector 의미 검색과 Elasticsearch 키워드 검색을 RRF로 병합")
    box(slide, 0.75, 1.45, 2.0, 0.65, "사용자 질문", PANEL_DARK, LINE, 12)
    box(slide, 3.45, 1.05, 2.8, 0.85, "pgvector\n의미 기반 검색", PURPLE, PURPLE, 12)
    box(slide, 3.45, 2.35, 2.8, 0.85, "Elasticsearch\nBM25 전문검색", YELLOW, YELLOW, 12)
    box(slide, 7.1, 1.7, 2.35, 0.85, "Application\nRetriever", GREEN, GREEN, 12)
    box(slide, 10.25, 1.7, 2.15, 0.85, "Ollama qwen3\n코드 기반 답변", BLUE, BLUE, 11)
    arrow(slide, 2.75, 1.78, 3.45, 1.48, PURPLE)
    arrow(slide, 2.75, 1.78, 3.45, 2.78, YELLOW)
    arrow(slide, 6.25, 1.48, 7.1, 2.0, PURPLE)
    arrow(slide, 6.25, 2.78, 7.1, 2.12, YELLOW)
    arrow(slide, 9.45, 2.12, 10.25, 2.12, BLUE)
    bullet_panel(slide, 1.0, 4.15, 3.45, 1.55, "pgvector", ["자연어 의미 유사도", "코드 맥락 검색", "가중치 1.0"], PURPLE)
    bullet_panel(slide, 4.95, 4.15, 3.45, 1.55, "Elasticsearch", ["함수명/파일명", "에러 문자열/식별자", "가중치 1.25"], YELLOW)
    bullet_panel(slide, 8.9, 4.15, 3.45, 1.55, "RRF Merge", ["중복 제거", "권한 필터", "Top-K 컨텍스트"], GREEN)
    footer(slide, 5)

    # 6. 업로드 흐름
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "파일 / 프로젝트 업로드 흐름")
    items = [("업로드", BLUE), ("언어 감지", PANEL), ("Tree-sitter\n파싱", PANEL), ("코드 청킹", PANEL), ("Ollama\n임베딩", PURPLE)]
    x = 0.55
    for i, (txt, color) in enumerate(items):
        box(slide, x, 1.65, 1.85, 0.78, txt, color, color if color != PANEL else LINE, 11)
        if i < len(items) - 1:
            arrow(slide, x + 1.85, 2.04, x + 2.25, 2.04)
        x += 2.25
    box(slide, 2.0, 4.25, 3.3, 0.95, "PostgreSQL / pgvector\ncode_files · code_chunks · embeddings", PURPLE, PURPLE, 11)
    box(slide, 7.35, 4.25, 3.3, 0.95, "Elasticsearch\ncodereview-code-chunks", YELLOW, YELLOW, 11)
    arrow(slide, 9.75, 2.43, 3.65, 4.25, PURPLE)
    arrow(slide, 9.75, 2.43, 9.0, 4.25, YELLOW)
    text_box(slide, 1.0, 6.25, 11.3, 0.35, "같은 코드 청크가 pgvector에는 벡터로, Elasticsearch에는 전문검색 문서로 저장됩니다.", 13, MUTED, False, PP_ALIGN.CENTER)
    footer(slide, 6)

    # 7. RAG / 이미지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "RAG Q&A와 이미지 기반 분석")
    bullet_panel(slide, 0.8, 1.35, 5.55, 4.9, "프로젝트 RAG Q&A", ["프로젝트 ZIP 업로드", "청킹 / 임베딩 / 색인", "Hybrid Retrieval", "Context Builder", "답변 + RAG 참조"], GREEN)
    bullet_panel(slide, 7.0, 1.35, 5.55, 4.9, "이미지 기반 분석", ["이미지 + 질문 입력", "llava 이미지 설명 생성", "설명을 검색 질의로 사용", "관련 코드 청크 검색", "qwen3가 코드 기반 답변"], PURPLE)
    footer(slide, 7)

    # 8. 비동기 리뷰
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "비동기 코드 리뷰")
    box(slide, 0.85, 2.0, 1.85, 0.75, "Review\nRequest", BLUE, BLUE, 11)
    box(slide, 3.15, 2.0, 1.85, 0.75, "FastAPI\nReview API", GREEN, GREEN, 11)
    box(slide, 5.45, 2.0, 1.85, 0.75, "Celery\nQueue", ORANGE, ORANGE, 11)
    box(slide, 7.75, 2.0, 1.85, 0.75, "Worker\n+ qwen3", PURPLE, PURPLE, 11)
    box(slide, 10.05, 2.0, 1.85, 0.75, "Review\nResult", YELLOW, YELLOW, 11)
    for x in [2.7, 5.0, 7.3, 9.6]:
        arrow(slide, x, 2.38, x + 0.45, 2.38)
    bullet_panel(slide, 2.0, 4.0, 9.2, 1.65, "상태 조회", ["pending", "processing", "completed", "failed"], CYAN)
    text_box(slide, 1.0, 6.25, 11.3, 0.35, "긴 LLM 리뷰 작업은 HTTP 요청 안에서 직접 처리하지 않고 Celery Worker가 처리합니다.", 13, MUTED, False, PP_ALIGN.CENTER)
    footer(slide, 8)

    # 9. 로그와 관측성
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "로그와 관측성")
    box(slide, 0.9, 1.4, 2.35, 0.8, "FastAPI\nMiddleware/Services", GREEN, GREEN, 11)
    box(slide, 4.25, 1.4, 2.35, 0.8, "Elasticsearch\nLogs", YELLOW, YELLOW, 11)
    box(slide, 7.6, 1.4, 2.35, 0.8, "Kibana\nData View", YELLOW, YELLOW, 11)
    box(slide, 10.7, 1.4, 1.75, 0.8, "Admin\nView", BLUE, BLUE, 11)
    arrow(slide, 3.25, 1.8, 4.25, 1.8)
    arrow(slide, 6.6, 1.8, 7.6, 1.8)
    arrow(slide, 9.95, 1.8, 10.7, 1.8)
    bullet_panel(slide, 1.0, 3.25, 3.4, 2.3, "수집 이벤트", ["API 요청", "RAG 검색", "LLM 호출", "웹 검색", "장애 이벤트"], GREEN)
    bullet_panel(slide, 5.0, 3.25, 3.4, 2.3, "인덱스", ["codereview-api-request-*", "codereview-rag-search-*", "codereview-llm-call-*", "codereview-web-search-*"], YELLOW)
    bullet_panel(slide, 9.0, 3.25, 3.4, 2.3, "확인 항목", ["latency", "failure rate", "empty RAG results", "token usage"], BLUE)
    footer(slide, 9)

    # 10. 포트와 요약
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "주요 포트와 요약")
    box(slide, 0.9, 1.35, 5.2, 3.55, "Next.js                  3000\nFastAPI                  8000\nOllama                   11434\nPostgreSQL + pgvector    5433 → 5432\nRedis                    6379\nElasticsearch            9200\nKibana                   5601", PANEL_DARK, LINE, 13, False, PP_ALIGN.LEFT)
    bullet_panel(slide, 7.0, 1.35, 5.2, 3.55, "요약", ["로컬 모델 기반 개발 보조 시스템", "pgvector + Elasticsearch Hybrid RAG", "SSE 스트리밍 응답", "Celery 비동기 리뷰", "Kibana 운영 관측성"], BLUE)
    text_box(slide, 1.0, 5.9, 11.3, 0.5, "핵심: 코드 컨텍스트 검색 품질, 로컬 LLM 답변, 운영 로그 관측성을 함께 갖춘 인트라넷형 AI 개발 도구", 14, TEXT, True, PP_ALIGN.CENTER)
    footer(slide, 10)

    return prs


if __name__ == "__main__":
    presentation = build_presentation()
    presentation.save(OUT)
    print(OUT)
